from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Cm, Inches, Pt
import numpy as np
from math import isnan
import six
import pandas as pd
import os

os.chdir(r"Z:\ResearchHome\Groups\millergrp\projects\CrisprDesigns\common\recent_designs\hAKT1-PMH-SORTTEST")

if not os.path.exists(r"Z:\ResearchHome\Groups\millergrp\projects\CrisprDesigns\common\recent_designs\hAKT1-PMH-SORTTEST\CRISPR_summary_w_primers_NGS.txt"):
    print("File not found")
    exit()
else:
    print("File found")


BIN_DIR = os.path.dirname(os.path.realpath(__file__))

print(BIN_DIR)


FINAL_SUMMARY = r"Z:\ResearchHome\Groups\millergrp\projects\CrisprDesigns\common\recent_designs\hAKT1-PMH-SORTTEST\CRISPR_summary_w_primers_NGS.txt"


DESIGNS_DIR = r"Z:\ResearchHome\Groups\millergrp\projects\CrisprDesigns\common\recent_designs"


def _do_formatting(value, format_str):
    """Format value according to format_str, and deal
    sensibly with format_str if it is missing or invalid.
    """
    if format_str == "":
        if type(value) in six.integer_types:
            format_str = ","
        elif type(value) is float:
            format_str = "f"
        elif type(value) is str:
            format_str = "s"
    elif format_str[0] == ".":
        if format_str.endswith("R"):
            if type(value) in six.integer_types:
                value = round_to_n(value, int(format_str[1]))
                format_str = ","
        if not format_str.endswith("G"):
            format_str = format_str + "G"
    try:
        value = format(value, format_str)
    except:
        value = format(value, "")

    return value

def process_position_parameter(param):
    """Process positioning parameters (left, top, width, height) given to
    df_to_table.
    If an integer, returns the right instance of the Cm class to allow it to be
    treated as cm. If missing, then default to 4cm. Otherwise, pass through
    whatever it gets.
    """
    if param is None:
        return Cm(4)
    elif type(param) is int:
        return Cm(param)
    else:
        return param

def df_to_table(
    slide,
    df,
    picked,
    left=Inches(1.07),#using non targeted sequences
    top=Inches(1.5),
    width=Inches(5.5),
    height=None,
    colnames=None,
    col_formatters=None,
    rounding=None,
    name="CRISPR_summary"):
    """Converts a Pandas DataFrame to a PowerPoint table on the given
    Slide of a PowerPoint presentation. Look for complete instructions at 
    https://github.com/robintw/PandasToPowerpoint/blob/master/pd2ppt/pd2ppt.py
    """
    
    #if targeted sequence is used, shift table left to keep everything aligned
    if len(df.axes[1]) == 7:
        left = Inches(0.68)
    
    left = process_position_parameter(left)
    top = process_position_parameter(top)
    width = process_position_parameter(width)
    height = process_position_parameter(height)
    

    rows, cols = df.head(8).shape
    shp = slide.shapes.add_table(rows + 1, cols, left, top, width, height)

    if colnames is None:
        colnames = list(df.columns)
    
    # Insert the column names
    for col_index, col_name in enumerate(colnames):
        shp.table.cell(0, col_index).text = col_name
        cell = shp.table.cell(0, col_index)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(9)

    m = df.values

    for row in range(rows):
        for col in range(cols):
            val = m[row, col]

            if col_formatters is None:
                text = str(val)
            else:
                text = _do_formatting(val, col_formatters[col])

            shp.table.cell(row + 1, col).text = text
            cell = shp.table.cell(row + 1, col)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(9)
            para.font.name = "Courier New"
            if row in picked:
                para.font.color.rgb = RGBColor(255, 0, 0)
            shp.table.cell(
                row, col
            ).text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            cell.text_frame.word_wrap = False
    shp.table.columns[0].width = Inches(1.95)
    shp.table.columns[1].width = Inches(2.25)
    
    ##
    try:
        shp.table.columns[6].width = Inches(1.3)
    except:
        pass
        # print("no column 8")

    if name is not None:
        shp.name = name

    return shp
    
def CAGE_slides(gene, species, project_dir, edit):
    if species == "hamster":
        s_gene = 'cho'
    else:
        s_gene = species[0] + gene
        
    presentation_filename = f"{project_dir}_{edit}_Summary.pptx"
    final_path = os.path.join(project_dir, presentation_filename)
    
    #input(final_path)
    
    prs = Presentation(os.path.join(BIN_DIR, "empty.pptx"))
    # slide 1
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    shapes.title.text = f"{s_gene} {edit} Summary"

    left = Inches(6)
    top = Inches(1.5)
    width = Inches(0.3)
    height = Inches(0.75)

    shape = shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 0, 0)
    line = shape.line
    line.color.rgb = RGBColor(255, 0, 0)

    # slide 2
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    shapes.title.text = f"{s_gene} gRNA(s)"
    top = width = height = Inches(1.5)
    left = Inches(2.75)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "gRNA sequences indicated with gray arrows"

    # slide 3
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Off Target Analysis - " + s_gene
    left = Inches(0.5) 
    top = Inches(5.25)  
    slide3_txtbox = slide.shapes.add_textbox(left, top, width, height)
    s3_txtbox_txt = slide3_txtbox.text_frame.paragraphs[0]
    s3_txtbox_txt.font.size = Pt(12)
    s3_txtbox_txt.text = """
    Long_0: the number sites in the genome that are an exact match to the full gRNA site, including the target site
    Long_1: the number sites in the genome that contain up to 1bp of mismatch to the full gRNA site, including the target site
    Long_2: the number sites in the genome that contain up to 2bp of mismatches to the full gRNA site, including the target site
    Long_3: the number sites in the genome that contain up to 3bp of mismatches to the full gRNA site, including the target site
    """
    
    #pic = slide.shapes.add_picture(os.path.join(BIN_DIR, "off_target.jpg"), left, top,) #* removed 11/6/23 PMH

    try:
        input_df = pd.read_csv(FINAL_SUMMARY, sep="\t")
                
        #TODO picked goes to the top
        #picked df concated to top of non picked.
        #sort each individually and merge
        
        
        picked_df = input_df[input_df["Picked"] == True].copy()
        
        secondary_df = input_df[input_df["Picked"] == False].copy()
        
        picked_df.sort_values(by=['long_0','long_1','long_2','long_3','Distance_from_BP'], inplace=True, ascending=True)
        picked_df.reset_index(inplace=True,drop=True)
        secondary_df.sort_values(by=['long_0','long_1','long_2','long_3'], inplace=True, ascending=True)




        sorted_df = pd.concat([picked_df, secondary_df])
        
        sorted_df.reset_index(inplace=True,drop=True)
        
        
        sorted_df.to_csv("sorted_table.txt", sep="\t", index=False)
        
        picked = picked_df.index.tolist()
        
        list_df = sorted_df.loc[picked, "Name"]
        # gRNA_list = df.loc[picked,'Name'].to_string(index = False)

        if isnan(sorted_df.loc[0, "Distance_from_BP"]):
            df1 = sorted_df.iloc[:, 0:6]
        else:
            #targeted sequence
            df1 = sorted_df.iloc[:, np.r_[0:6, 7]]
        df1 = df1.fillna("NA")
        
        df1.to_csv("df1.txt", sep="\t", index=False)
        
        table = df_to_table(slide, df1, picked)
    except:
        print("Trouble making table in presentation. Please make manually.")
        
        
        
    # slide 4
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Summary"

    tf = body_shape.text_frame
    tf.text = "The following gRNA(s) are recommended based on off-target profile and the parameters of your project:"

    try:
        for i in list_df:
            p = tf.add_paragraph()
            p.text = i
            p.level = 1
    except:
        p = tf.add_paragraph()
        p.text = "add your gRNA names here"
        p.level = 1

    # ppt created
    try:
        prs.save(final_path)
    except Exception as e:
        print("Error saving ppt")
        print(e)    
        
project_dir = r'Z:\ResearchHome\Groups\millergrp\projects\CrisprDesigns\common\recent_designs\hAKT1-PMH-SORTTEST\\'
        
ppt_path = CAGE_slides('CHCHD10-test', 'human', project_dir, 'KO')