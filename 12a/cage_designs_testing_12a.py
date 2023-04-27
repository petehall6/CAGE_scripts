#!/usr/bin/env python
import argparse
import collections
import glob
import itertools
import logging
import os
import pickle
import re
import shutil
import time
from math import isnan

from pathlib import PurePosixPath, PureWindowsPath

import numpy as np
import pandas as pd
import regex
import six
from Bio import Entrez, SeqFeature, SeqIO, SeqUtils

from Bio.Seq import Seq
from Bio.SeqRecord import SeqRecord
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Cm, Inches, Pt

import crisprSummarySRMConverter
from check_existing_projects import check_projects
from off_target_all_RDM_v4_12a import off_target_all, off_target_all_primers

# os.system("module load gcc/4.8.5")

# These variables are full-on global because they are used by either tests
# or interactive_crispr.py
BIN_DIR = os.path.dirname(os.path.realpath(__file__))
DESIGNS_DIR = "/research_jude/rgs01_jude/groups/millergrp/projects/CrisprDesigns/common/recent_designs"
Z_DESIGNS_DIR = PureWindowsPath(
    r"Z:\ResearchHome\Groups\millergrp\home\common\CORE PROJECTS"
)
Z_DESIGNS_FOR_CLUSTER = PurePosixPath(
    "/research_jude/rgs01_jude/groups/millergrp/home/common/CORE\ PROJECTS"
)

GENE_GENBANK = "gene.gbk"
# GENBANK_LENGTH_EXTENSION = 5000
PROJECT_PARAMS = "params"
LOG_FILE = "log.txt"

OFF_TARGET_BOWTIE = "CRISPR_long_0mm.bowtie.formatted"

# Column names for CRISPR_summary* files
SNP_DATAFRAME_COLUMNS = (
    "SNP",
    "SNP_count",
)  # used in create_snp_dataframe() and crispr_design()
CUSTOMER_COLUMN = "Customer"
SEQUENCE_COLUMN = "gRNA"

SUPPORTED_SPECIES = ["human", "mouse", "zebrafish", "rat", "hamster", "horse", "monkey"]
SPECIES_TO_ABBRV = {
    "human": "h",
    "mouse": "m",
    "zebrafish": "z",
    "rat": "r",
    "hamster": "ch",
    "horse": "e",
    "monkey": "gm",
}
ABBRV_TO_SPECIES = {
    "h": "HUMAN",
    "m": "mouse",
    "z": "zebrafish",
    "r": "rat",
    "ch": "hamster",
    "e": "horse",
    "gm": "monkey",
}

VECTOR_CHOICES = [
    "SYN",
    "SM115",
    "SM102",
    "SM104",
    "SM387",
    "SM388",
    "SP16",
    "SP82",
    "SM438",
    "PC323",
    "SNP23",
    "SNP21",
    "SNP22",
    "SNP36",
    "SNP329",
    "SNP27",
    "SNP38",
    "SNP28",
    "SNP112",
    "SP76",
    "SP78",
    "PC291"
]
# Outputed file names
OFF_TARGET_SUMMARY = "CRISPR_off_target_v3.summary"
CRISPR_SUMMARY = "CRISPR_summary.txt"
FINAL_SUMMARY = "CRISPR_summary_w_primers_NGS.txt"
AMPLICON_REPORT = "PCR_assay_report.txt"

BASIC_SUMMARY_SORT_KEY = ["long_0", "long_1", "long_2", "short_0"]

ENSEMBL_RELEASE_ARCHIVE = "mar2015.archive.ensembl.org"  # orignial archive used
# ENSEMBL_RELEASE_ARCHIVE = "apr2020.archive.ensembl.org"   #hangs
ENSEMBL_RELEASE_VERSION = "75"  # not actually used, just for documentation

USER_EMAIL = "aschriefer@path.wustl.edu"
Entrez.email = USER_EMAIL

# Thrown by SeqTargeted.__init__
class InvalidBasesError(Exception):
    pass


class MismatchedTargetingCharError(Exception):
    pass


# Thrown by primer_design() and GuideRNA.find_target()
class InvalidFuzzyStringError(Exception):
    pass

class txtcolors:
    '''
    Black   30
    Red     31
    Green   32
    Yellow  33
    Blue    34
    Magenta 35
    Cyan    36
    White   37
    Bright Black    90
    Bright Red      91
    Bright Green    92
    Bright Yellow   93
    Bright Blue     94
    Bright Magenta  95
    Bright Cyan     96
    Bright White    97
    '''
    #ANSI escape sequence to set back ground color.  Replace int after [
    GOOD = '\033[92m'
    WARNING = '\033[91m'
    HIGHLIGHT = '\033[95m'
    #Clears color back to default
    END = '\033[0m'
class Amplicon(object):
    def __init__(self, Seq_target, penalty, fwd_primer, rev_primer, start, end):
        """
        self.seq is sliced with end+1 to convert to python slicing
        from raw primer3 base coordinates
        """
        self.fwd_primer = fwd_primer
        self.rev_primer = rev_primer
        self.penalty = float(penalty)
        self.start = int(start)
        self.end = int(end)
        self.seq = Seq_target[self.start : (self.end + 1)]

    def fragment(self, guide_target_list):
        cut_points = [x.cut for x in guide_target_list]
        
        sorted_cut_points = sorted(cut_points)
        
        assert sorted_cut_points[-1] < len(self.seq)

        break_points = [0] + sorted_cut_points + [len(self.seq)]
        frag_lengths = [c2 - c1 for c1, c2 in zip(break_points[:-1], break_points[1:])]

        nonzero_frag = [f for f in frag_lengths if f > 0]
        return nonzero_frag

    def __str__(self):
        return " ".join([self.fwd_primer, self.rev_primer])


class PrimerDesigner(object):
    """
    Finds a region in seq_record which contains all guide_rnas in guide_rna_list
    extends the region by MAX_PRODUCT_SIZE in either direction and passes the region
    to primer3 for primer design.  This is controlled by the only public method
    create_ampliconSet which returns an AmpliconSet object.  The AmpliconSet object
    is responsible for validating the output amplicons from primer3
    """

    # PRIMER_DESIGN_BIN = '/home/comp/gtac/gtac/software/primer3-2.3.6/src/primer3_core'
    PRIMER_DESIGN_BIN = "/research_jude/rgs01_jude/groups/millergrp/projects/CrisprDesigns/common/local/primer3-2.3.6/src/primer3_core"
    # contains global settings for all primer3 runs
    PRIMER_GLOBAL_SETTINGS = os.path.join(BIN_DIR, "primer3_global_settings")
    # input file for primer3, contains target seq and specific settings
    PRIMER_TARGET = "primer3_target"
    PRIMER_OUTPUT = "primer_set.txt"

    MAX_PRODUCT_SIZE = 1000  # largest acceptable pcr product
    GRNA_BUFFER_SIZE = 110  # adds bases to either end of primer3 included region
    PRODUCT_SIZE_RANGE = "120-250 251-275 276-300 301-400"
    PRIMER_NUM_RETURN = "10"

    def __init__(
        self, guide_rna_list, seq_record, grna_picking_seq, fuzzy_match_control=None
    ):
        """
        Assumes that all GuideRNAs in guide_rna_list are perfectly contained
        within seq_picked.  seq_record is a Bio.SeqRecord, grna_picking_seq is a string or Bio.Seq
        """
        self.guide_rna_list = guide_rna_list
        # location of grna wrt the genbank sequence
        # simple 2-tuples NOT GuideTargets
        self.fuzzy = fuzzy_match_control
        self.seq_record = seq_record
        self.grna_picking_seq = grna_picking_seq

    def create_ampliconSet(self, args):
        """
        Makes primers in the current working directory
        """
        self.relative_grna_start = self._find_grna_region()

        grna_region, target_seq = self._get_target_seq()
        # print(f"208: {grna_region}") # Region will be padded by GRNA_BUFFER_SIZE
        grna_start, grna_end = grna_region
        grna_length = grna_end - grna_start

        assert (
            grna_length <= self.MAX_PRODUCT_SIZE
        ), "total region encompassing all gRNA is larger than MAX_PRODUCT_SIZE"

        with open(self.PRIMER_TARGET, "w") as fh:
            fh.write(
                "SEQUENCE_ID=%s\n" % ";".join([x.name for x in self.guide_rna_list])
            )
            fh.write("SEQUENCE_TEMPLATE=%s\n" % target_seq)
            fh.write("SEQUENCE_TARGET=%s,%s\n" % (grna_start, grna_length))
            fh.write("PRIMER_PRODUCT_SIZE_RANGE=%s\n" % self.PRODUCT_SIZE_RANGE)
            fh.write("PRIMER_NUM_RETURN=%s\n" % self.PRIMER_NUM_RETURN)
            fh.write("=")

        primer_call = "{cmd} -echo_settings_file -p3_settings_file {setting} -output {out} {infile}".format(
            cmd=self.PRIMER_DESIGN_BIN,
            setting=self.PRIMER_GLOBAL_SETTINGS,
            out=self.PRIMER_OUTPUT,
            infile=self.PRIMER_TARGET,
        )
        # print(f"232: Primer call\n{primer_call}\n")
        os.system(primer_call)

        amplicon_list = self._parse_primer_output(target_seq, self.PRIMER_OUTPUT, grna_start, grna_end, args)
        # print(amplicon_list)
        return AmpliconSet(amplicon_list, self.guide_rna_list, self.fuzzy)

    def _get_target_seq(self):
        """
        Outputs the sequence that will be passed to primer design program ie the target
        Also gives the start and end of this target that contain the grnas ie the
        regions that MUST be covered by the primers
        """
        rel_start_sort = sorted(self.relative_grna_start, key=lambda x: x[0])
        grna_region_start = rel_start_sort[0][0]
        grna_region_end = rel_start_sort[-1][1]

        target_start, target_end = self._extend_target(
            len(self.seq_record),
            self.MAX_PRODUCT_SIZE,
            grna_region_start,
            grna_region_end,
        )

        primer_template_seq = str(self.seq_record.seq[target_start : target_end + 1])

        # Represents coordinates of grna region in terms of the target
        relative_grna_start = grna_region_start - target_start
        relative_grna_end = grna_region_end - target_start

        # Add extra room that must be covered by primers in order to ensure
        # grna cut gives fragment with at least self.GRNA_BUFFER_SIZE length
        buffered_grna_start, buffered_grna_end = self._extend_target(
            len(primer_template_seq),
            self.GRNA_BUFFER_SIZE,
            relative_grna_start,
            relative_grna_end,
        )

        return ((buffered_grna_start, buffered_grna_end), primer_template_seq)

    # The start and ends are in terms of element numbers not slicing indicies
    def _extend_target(self, N, extend, g_start, g_end):
        assert extend >= 0, "Cannot extend with negative number %s" % extend
        t_start = max(0, g_start - extend)
        t_end = min(N, g_end + extend)

        return (t_start, t_end)

    def _find_grna_region(self):
        """
        First align sequence that was used to find gRNAs in the crispr_design routine
        Then aligns guide_rnas to that sequence and shifting over to find absolute
        position of the guide_rnas
        """
        # create a dummy GuideRNA object to use find_target
        # This is checking the entire gbk for your input grna picking seq
        # Should show up exactly once. No more, no less.
        # One shalt be the number of occurrences, and the number of occurrences should be one.
        # print("picking target")
        picking_target = GuideRNA("picked", self.grna_picking_seq).find_target(self.seq_record.seq, fuzzy=self.fuzzy)
      
        # print(self.seq_record.seq)
        # print(picking_target)
        # print(self.grna_picking_seq)

        # Dear god I hope this never happens --SP
        
        assert len(picking_target) == 1, (
            "GRNA picking target has %s matches with fuzzy %s"
            % (len(picking_target), self.fuzzy)
        )

        logging.info("Designing primers with fuzzy match control : %s" % self.fuzzy)

        picking_target_start = picking_target[0].start
        relative_grna_start = []
        for guide_rna in self.guide_rna_list:
            guideTarget = guide_rna.find_target(self.grna_picking_seq)

            # warn user in log that the gRNA has multiple/no targets in the
            # passed genbank file

            # I need to think about this whole multiple guideRNA binding sites deal
            N = len(guideTarget)
            if N != 1:
                logging.warn(
                    "gRNA {gid} has {n} targets in seq_record".format(
                        gid=guide_rna.name, n=N
                    )
                )

            assert N > 0, (
                "Guide rna %s has no matches in grna_picking_seq" % guide_rna.name
            )

            grna_target = [
                (picking_target_start + t.start, picking_target_start + t.end)
                for t in guideTarget
            ]

            relative_grna_start.extend(grna_target)

        return relative_grna_start

    def _parse_primer_output(
        self, target_seq, primer_output_fp, grna_start, grna_end, args
    ):
        """
        Assumes that there are no spaces between field name, =, and entry
        in the primer_output_fp
        """
        primer_records = dict(
            line.strip().split("=") for line in open(self.PRIMER_OUTPUT, "r")
        )
        amplicon_list = []
        iprimer = 0
        while "PRIMER_LEFT_{0}_SEQUENCE".format(iprimer) in primer_records:
            amplicon = Amplicon(
                target_seq,
                primer_records["PRIMER_PAIR_{0}_PENALTY".format(iprimer)],
                primer_records["PRIMER_LEFT_{0}_SEQUENCE".format(iprimer)],
                primer_records["PRIMER_RIGHT_{0}_SEQUENCE".format(iprimer)],
                primer_records["PRIMER_LEFT_{0}".format(iprimer)].split(",")[0],
                primer_records["PRIMER_RIGHT_{0}".format(iprimer)].split(",")[0],
            )

            amplicon_list.append(amplicon)
            iprimer += 1

        logging.info("Found %s pcr assay primers" % len(amplicon_list))

        continue_w_primers = ""
        # print("amplicon list")
        # for amplicon in amplicon_list:
        # print(amplicon)
        # print("amplicon list length",len(amplicon_list))
        ota_run = False

        if len(amplicon_list) > 0:
            primer_dict = {}
            i = 0
            for item in amplicon_list:
                fwd = item.fwd_primer
                rev = item.rev_primer
                fwd_dict = {f"forward primer {i}": fwd}
                rev_dict = {f"reverse primer {i}": rev}
                primer_dict.update(fwd_dict)
                primer_dict.update(rev_dict)
                i = i + 1

            check_primer_off_targets_dict(primer_dict, args.project_dir)
            ota_run = True

        # continue_w_primers = input("Continue with listed primers? (y/n)")

        # if continue_w_primers.upper() == "N" or len(amplicon_list) <= 0:
        if len(amplicon_list) <= 0:
            print("\n* * *   No good pcr primers found!   * * *\n")

        try:
            choice = input(
                "Would you like to manually enter primer sequences? (y/n)\n  -> "
            )
            if choice.upper() == "Y":
                ota_run = False
                print()
                forward_primer = input("    Forward primer -> ").strip().upper()
                reverse_primer = input("    Reverse primer -> ").strip().upper()
                reverse_primer = str(Seq(reverse_primer))

                forward_location = str(self.seq_record.seq).find(forward_primer)
                forward_location_end = forward_location + len(forward_primer)
                # print(
                #     f"    \nforward primer location: ({forward_location}, {forward_location_end})"
                # )

                reverse_location = str(self.seq_record.seq).find(
                    str(Seq(reverse_primer).reverse_complement())
                )
                reverse_location_end = reverse_location + len(reverse_primer)
                # print(
                #     f"    reverse primer location: ({reverse_location}, {reverse_location_end})"
                # )

                grna_se_list = self.relative_grna_start
                se_tuple = grna_se_list[0]
                grna_start, grna_end = se_tuple
                # print(f"    gRNA location: ({grna_start}, {grna_end})")

                dif1 = abs(grna_start - forward_location_end)
                dif2 = abs(reverse_location - grna_end)

                if (
                    dif1 < 50
                    or dif2 < 50
                    or grna_end >= reverse_location
                    or grna_start <= forward_location_end
                ):
                    print(
                        "\n* gRNA and primers are likely not compatible (overlap, close proximity, or gRNA not in amplicon) *\n"
                    )
                else:
                    print("primers look good!")
                    choice = "N"

                amplicon_seq = self.seq_record.seq[
                    forward_location : reverse_location_end + 1
                ]
                # print(f"amplicon_seq:\n  z{amplicon_seq}")

                amplicon = Amplicon(
                    self.seq_record.seq,  # Was target_seq
                    0,
                    forward_primer,
                    reverse_primer,
                    # target_seq.upper().find(forward_primer.upper()),
                    forward_location,
                    reverse_location,
                )
                amplicon_list.insert(0, amplicon)

            if len(amplicon_list) > 0 and ota_run == False:
                primer_dict = {}
                i = 0
                for item in amplicon_list:
                    fwd = item.fwd_primer
                    rev = item.rev_primer
                    fwd_dict = {f"forward primer {i}": fwd}
                    rev_dict = {f"reverse primer {i}": rev}
                    primer_dict.update(fwd_dict)
                    primer_dict.update(rev_dict)
                    i = i + 1

                check_primer_off_targets_dict(primer_dict, args.project_dir)
                ota_run = True

        except EOFError:
            return None
        # assert len(amplicon_list) > 0, 'Error! No pcr primers found!'
        return amplicon_list


class AmpliconSet(object):
    MIN_FRAGMENT_DELTA = 1
    MIN_FRAG_LENGTH = 1

    def __init__(self, amplicon_list, guide_rna_list, fuzzy_match_control):
        self.amplicon_list = amplicon_list
        self.guide_rna_list = guide_rna_list
        self.fuzzy = fuzzy_match_control

    def choose_best_amplicon(self, name=None):
        """
        Returns a tuple (amplicon object, dataframe)
        Assumes that self.amplicon_list is ordered from best to worst
        Pass name to set the amplicon name in the dataframe
        otherwise name is given by self._make_amplicon_name
        """
        # print("(492) choose_best")
        # print(self.amplicon_list)
        for amplicon in self.amplicon_list:
            # print(f"Amplicon:    {amplicon}")
            is_amplicon_valid = self._is_amplicon_valid(amplicon)

            if is_amplicon_valid:  # Only returns for first valid amplicon

                if round(SeqUtils.GC(amplicon.seq),3) > 70:
                    print(f"{txtcolors.WARNING}")
                    print("*"*50)
                    print("*"*50)
                    print("HIGH GC CONTENT WARNING:")
                    print(f"Amplicon GC content is: {round(SeqUtils.GC(amplicon.seq),3)}")
                    print("Please let Batch Team Know")
                    print("*"*50)
                    print("*"*50)
                    input("Press any key to continue")
                    print(f"{txtcolors.END}")
                else:
                        # print("valid")
                        print("\n")
                        print(f"{txtcolors.HIGHLIGHT}")
                        print("*"*50)
                        print("*"*50)
                        print(f"Amplicon GC content is: {round(SeqUtils.GC(amplicon.seq),3)}")
                        print("*"*50)
                        print("*"*50)
                        print("\n")
                        input("Press any key to continue")
                        print(f"{txtcolors.END}")
                        amplicon_df = self._create_amplicon_dataframe(amplicon, name=name)
                        amplicon_df.drop("Penalty", axis=1, inplace=True)
                return (amplicon, amplicon_df)

        # should never reach this point if there exists a valid amplicon
        assert False, "No valid primer pairs were found!"

    def write_report_for_N_best_amplicons(self, N, output_fp, name=None):
        """
        Pass name to set all the amplicon names to name
        """
        numChosen = 0
        # print(f"(513) amplicon_list has {len(self.amplicon_list)} amplicons.")

        for amplicon in self.amplicon_list:
            # print(f"Is {amplicon} valid?")
            is_amplicon_valid = self._is_amplicon_valid(amplicon)

            if is_amplicon_valid:
                amplicon_df = self._create_amplicon_dataframe(amplicon, name=name)
                self._write_amplicon_report_file(amplicon_df, output_fp)

                numChosen += 1

            if numChosen >= N:
                break

        logging.info(f"Requested {N} valid primer sets; found {numChosen}")
        return output_fp

    def _are_fragments_valid(self, grna_to_frag_lengths):
        # print(f"Are fragments valid?  ({grna_to_frag_lengths})")
        fragment_bool_list = []
        for fragment_lengths in grna_to_frag_lengths.values():
            frag_combos = itertools.combinations(fragment_lengths, 2)
            frag_length_diff = [abs(x[0] - x[1]) for x in frag_combos]

            fragment_bool_list.append(
                min(frag_length_diff) >= self.MIN_FRAGMENT_DELTA
                and min(fragment_lengths) >= self.MIN_FRAG_LENGTH
            )
        # print("Are fragments valid?")
        return all(fragment_bool_list)

    def _are_primers_uniquely_binding(self, primers):
        # TO DO:
        Exception("not implemented yet")

    def _are_guide_rna_covered(self, amplicon, guide_rna_list):
        """
        Run time check of output of PrimerDesigner
        """
        # print("Are gRNAs Covered?")
        covered_bool_list = [
            len(g.find_target(amplicon.seq, self.fuzzy)) > 0 for g in guide_rna_list
        ]

        return all(covered_bool_list)

    def _is_amplicon_valid(self, amplicon):
        # print(f"Is amplicon valid?  {amplicon}")
        frag_lengths = self._make_fragment_lengths(amplicon)
        # Create a method _are_primers_unique thats blasts the primers
        # to also ensure a valid amplicon has unique primers
        fragval = self._are_fragments_valid(frag_lengths)
        gcov = self._are_guide_rna_covered(amplicon, self.guide_rna_list)
        is_amplicon_valid = fragval and gcov
        # print(f"Is amplicon valid? {fragval} and {gcov}")
        return is_amplicon_valid

    def _make_amplicon_id(self, amplicon):
        hex_string = hex(hash(str(amplicon.seq)))
        N = len(hex_string)
        return hex_string[max(0, N - 10) :]

    def _make_amplicon_name(self):
        name = self.guide_rna_list[0].name
        for guide_rna in self.guide_rna_list[1:]:
            name += "." + guide_rna.get_relative_id()

        return name

    def _make_fragment_lengths(self, amplicon):
        cut_indexes = {
            g.name: g.find_target(amplicon.seq, self.fuzzy) for g in self.guide_rna_list
        }


        frag_lengths = {
            guide_name: amplicon.fragment(targets)
            for guide_name, targets in cut_indexes.items()
        }
        return frag_lengths

    def _create_amplicon_dataframe(self, amplicon, name=None):
        frag_lengths = self._make_fragment_lengths(amplicon)
        frag_strings = {g: ":".join(map(str, f)) for g, f in frag_lengths.items()}

        df_index = [g.name for g in self.guide_rna_list]

        assert set(frag_strings.keys()) == set(df_index)

        amplicon_id = self._make_amplicon_id(amplicon)
        amplicon_name = self._make_amplicon_name() if name is None else name
        amplicon_df = pd.DataFrame(
            {
                "Amplicon_ID": amplicon_id,
                "Penalty": amplicon.penalty,
                "Amplicon_Length": len(amplicon.seq),
                "Fragments": frag_strings,
                "Fwd_Primer_Name": amplicon_name + ".F",
                "Fwd_Primer": amplicon.fwd_primer,
                "Rev_Primer_Name": amplicon_name + ".R",
                "Rev_Primer": amplicon.rev_primer,
                "DS.Fwd_Primer_Name": amplicon_name + ".DS.F",
                #'DS.Fwd_Primer' : 'CACTCTTTCCCTACACGACGCTCTTCCGATCT'+amplicon.fwd_primer,
                "DS.Fwd_Primer": "CTACACGACGCTCTTCCGATCT" + str(amplicon.fwd_primer).lower(),
                "DS.Rev_Primer_Name": amplicon_name + ".DS.R",
                #'DS.Rev_Primer' : 'GTGACTGGAGTTCAGACGTGTGCTCTTCCGATCT' + amplicon.rev_primer},
                "DS.Rev_Primer": "CAGACGTGTGCTCTTCCGATCT" + str(amplicon.rev_primer).lower(),
            },
            index=df_index,
        )

        col_order = [
            "Amplicon_ID",
            "Penalty",
            "Amplicon_Length",
            "Fragments",
            "DS.Fwd_Primer_Name",
            "DS.Fwd_Primer",
            "DS.Rev_Primer_Name",
            "DS.Rev_Primer",
            "Fwd_Primer_Name",
            "Fwd_Primer",
            "Rev_Primer_Name",
            "Rev_Primer",
        ]
        return amplicon_df.reindex(col_order, axis=1)

    def _write_amplicon_report_file(self, amplicon_df, output_fp):
        amplicon_df.insert(
            3, "gRNA", pd.Series(amplicon_df.index, index=amplicon_df.index)
        )
        if os.path.isfile(output_fp):
            amplicon_df.to_csv(output_fp, sep="\t", mode="a", header=False, index=False)
        else:
            amplicon_df.to_csv(output_fp, sep="\t", index=False)


class GenBankFeatureWriter(object):
    def __init__(self, gbk_file):
        self.genbank_record = SeqIO.read(open(gbk_file, "r"), format="genbank")

    def add_features(self, feature_tuples):
        """
        Takes an input a list of tuples of form (qualifier_dict, type, sequence)
        This method aligns the sequence to the genbank record and writes
        a feature at the aligned location(s)
        """
        for feature in feature_tuples:
            for target in self._align_seq(feature[2]):
                if target.strand == "1":
                    strand = 1
                elif target.strand == "2":
                    strand = -1
                else:
                    raise Exception("invalid strand %s" % target.strand)

                # I'm pretty sure that genbank files have start index 0 and
                # end index 1, at least thats how snapgene appears to be reading it
                location = SeqFeature.FeatureLocation(
                    target.start, target.end + 1, strand=strand
                )

                Seq_feature = SeqFeature.SeqFeature(
                    location=location, type=feature[1], qualifiers=feature[0]
                )
                self.genbank_record.features.append(Seq_feature)

        return None

    def write(self, output_fp):
        SeqIO.write([self.genbank_record], open(output_fp, "w"), "genbank")

    def _align_seq(self, seq):
        feature = GuideRNA("", seq)
        feature.LENGTH = len(seq)
        # print("_align_seq")
        return feature.find_target(self.genbank_record.seq)


# Return class of GuideRNA.find_target()
GuideTarget = collections.namedtuple("GuideTarget", ["start", "end", "cut", "strand"])
orthos_list = {}
ortho_options = []
orthos_df = pd.read_csv(
    "/research_jude/rgs01_jude/groups/millergrp/projects/CrisprDesigns/common/crispr_designs/ortholog12a.csv",
    usecols=["Ortholog", "LENGTH", "CUT_SITE"],
    low_memory=True,
    index_col=False,
)

orthos_list = orthos_df.to_dict("list")

class GuideRNA(object):
    """
    Represents the Cas9/Crispr guide rna.
    But really its a pretty abstract class that allows any sequence to be aligned
    to a larger sequence as long as you ignore the cut field of return GuideTarget.
    """


    def __init__(self, name, seq):
        self.name = name
        self.seq = seq if isinstance(seq, Seq) else Seq(seq)
    def get_length(self):
        
        df = read_crispr_summary(CRISPR_SUMMARY)
        
        
        ##get cas type for matching guide ID
        df_summary_picked = df['Picked']
        #OG df combines Name and Seq with 'Name' set as index
        #reset index to access index as regular column
        df_temp = df_summary_picked.reset_index()
        #expand 'Name' column into sub strings to parse
        df_temp[['CAGENum','Gene','Cas','GuideNum']] = df_temp.Name.str.split('.', expand=True)
        #drop redundent name column
        df_picked = df_temp.drop(['Name'], axis=1)
        #create casType df with only picked guides
        df_casType = df_picked.loc[df_picked['Picked'] == True]
        #use iloc since some df's will only have 1 row
        ortho = df_casType.iloc[0]['Cas']
        
        if ortho == "Cas9":
            LENGTH = 23
        elif ortho == "Cas12a":
            LENGTH = 25

    def find_target(self, Seq_ref, fuzzy=None):
        """
        Takes Bio.Seq.Seq object Seq_ref with the DNA alphabet
        Returns a list of GuideTarget objects with (0 = start, 1 = end, 2 = cut location, 3=strand)
        Strand is a string '1' for forward strand (relative to Seq_ref) and '2' for reverse
        Start and end are reversed (ie PAM sequence at start) in the case of reverse string matches
        Cut is the python slicing index that the cut will be made (ie in-between actual bases)
        """
        if (
            self.name != "picked" and self.name != ""
        ):  # Added "and" statement to check self.name != ""
            ortho = self.name.split(".")[2]
        else:
            ortho = self.name
        if fuzzy is None:
            fuzzy = ""
        elif not self.validate_fuzzy_match(fuzzy):
            raise InvalidFuzzyStringError(fuzzy)

        ###GET CAS TYPE
        ##Read grna crispr summary
        df = read_crispr_summary(CRISPR_SUMMARY)
        ##get cas type for matching guide ID
        df_summary_picked = df['Picked']
        #OG df combines Name and Seq with 'Name' set as index
        #reset index to access index as regular column
        df_temp = df_summary_picked.reset_index()
        #expand 'Name' column into sub strings to parse
        df_temp[['CAGENum','Gene','Cas','GuideNum']] = df_temp.Name.str.split('.', expand=True)
        #drop redundent name column
        df_picked = df_temp.drop(['Name'], axis=1)
        #create casType df with only picked guides
        df_casType = df_picked.loc[df_picked['Picked'] == True]
        #use iloc since some df's will only have 1 row
        ortho = df_casType.iloc[0]['Cas']
        
        fwd_targets = []
        rev_targets = []
        fwd_matches = self._recog_site_fuzzy_regex(str(self.seq), fuzzy).finditer(
            str(Seq_ref)
        )
        if ortho == "Cas9":
            CUT_SITE = int(orthos_df.loc[orthos_df["Ortholog"] == "Cas9", "CUT_SITE"].iloc[0])
            #LENGTH = int(orthos_df.loc[orthos_df["Ortholog"] == "Cas9", "LENGTH"].iloc[0])
            
            fwd_targets = [
                GuideTarget(m.start(), m.end() - 1, m.start() + CUT_SITE, "1")
                for m in fwd_matches
            ]
            
        elif ortho == "Cas12a":
            CUT_SITE = int(orthos_df.loc[orthos_df["Ortholog"] == "Cas12a", "CUT_SITE"].iloc[0])
            #LENGTH = int(orthos_df.loc[orthos_df["Ortholog"] == "Cas12a", "LENGTH"].iloc[0])
            fwd_targets = [
                GuideTarget(m.start(), m.end() - 1, m.start() + CUT_SITE, "1")
                for m in fwd_matches
            ]

        rev_matches = self._recog_site_fuzzy_regex(
            str(self.seq.reverse_complement()), fuzzy
        ).finditer(str(Seq_ref))
        if ortho == "Cas9":
            CUT_SITE = int(
                orthos_df.loc[orthos_df["Ortholog"] == "Cas9", "CUT_SITE"].iloc[0]
            )
            rev_targets = [
                GuideTarget(m.start(), m.end() - 1, m.end() - CUT_SITE, "2")
                for m in rev_matches
            ]
        elif ortho == "Cas12a":
            CUT_SITE = int(
                orthos_df.loc[orthos_df["Ortholog"] == "Cas12a", "CUT_SITE"].iloc[0]
            )
            rev_targets = [
                GuideTarget(m.start(), m.end() - 1, m.end() - CUT_SITE, "2")
                for m in rev_matches
            ]
        return fwd_targets + rev_targets

    def get_relative_id(self):
        return self.name.split(".")[-1]

    @classmethod
    def validate_fuzzy_match(cls, fuzzy):
        m = re.match(r"{(?:[ids]<=\d+,)*(?:[ids]<=\d+)}", fuzzy)
        # Match at least one, but more if separated by commas
        # Should be fine with re module instead of regex
        # i.e. "{i=0,d=0,s<=1}"
        # print(f"validate m is :   {m}")
        return bool(m)

    def _recog_site_fuzzy_regex(self, recog_site, fuzzy=""):

        notation_map = {
            "A": "[A]",
            "C": "[C]",
            "G": "[G]",
            "T": "[T]",
            "I": "[AT]",
            "W": "[AT]",
            "S": "[CG]",
            "M": "[AC]",
            "K": "[GT]",
            "R": "[AG]",
            "Y": "[CT]",
            "B": "[CGT]",
            "D": "[AGT]",
            "H": "[ACT]",
            "V": "[ACG]",
            "N": "[ACGT]",
        }

        # Regex pattern w/ each nucleotide (nt) as [] around acceptable options
        regexp = "".join([notation_map[nt] for nt in recog_site.upper()])
        regexp = f"(?:{regexp})" + f"{fuzzy}"
        # print(f"recog site regexp is :   {regexp}")
        return regex.compile(regexp)

    def __str__(self):
        return str(self.seq)


class UniqueGuideRNA(GuideRNA):
    """
    Represents a Cas9 guide rna with a unique binding location on the genome
    """

    def __init__(self, name, seq, chrom, start):
        self.chr = str(chrom)
        self.start = int(start)
        self.name = name
        if self.name != "picked":
            self.ortho = self.name.split(".")[2]
        else:
            self.ortho = self.name

        if self.ortho == "Cas9":
            self.LENGTH = int(
                orthos_df.loc[orthos_df["Ortholog"] == "Cas9", "LENGTH"].iloc[0]
            )
        elif self.ortho == "Cas12a":
            self.LENGTH = int(
                orthos_df.loc[orthos_df["Ortholog"] == "Cas12a", "LENGTH"].iloc[0]
            )
        self.end = self.start + self.LENGTH - 1

        super(UniqueGuideRNA, self).__init__(name, seq)

    def info(self):
        """
        Serializes the object to be passed into ensembl 
        querey R script
        """
        return ",".join([self.name, self.chr, str(self.start), str(self.end)])


class SeqTargeted(Seq):
    TARGET_PATTERN = re.compile(r"<([GATCRYWSMKHBVDN])>")

    def __init__(self, seq):
        """
        self.target_indexes is empty list if no targets are found
        Accepts <A> as targeting markup. If any loose < or > are found
        raises an error.
        Raises an error if multiple bases are targeted with one set of <>
        """

        self.raw_sequence = seq
        pre_clean_seq = re.sub(r"[\'\"]+", "", seq.upper())

        # Validate the passed sequence string
        clean_seq = self.TARGET_PATTERN.sub(lambda m: m.group(1), pre_clean_seq)
        if re.search(r"[^GATCRYWSMKHBVDN<>]", clean_seq):
            raise InvalidBasesError(clean_seq)

        elif re.search(r"[<>]", clean_seq):
            raise MismatchedTargetingCharError(clean_seq)

        # Find the targets and adjust indexes to account for target chars
        targets = [
            m.start() - 2 * i
            for i, m in enumerate(self.TARGET_PATTERN.finditer(pre_clean_seq))
        ]

        self.target_indexes = targets

        super(SeqTargeted, self).__init__(clean_seq)


def create_guideRna_list_from_summary(summary_fp, sequence_column):
    df = read_crispr_summary(summary_fp)
    return [GuideRNA(str(ind), df.loc[ind, sequence_column]) for ind in df.index]


def create_unique_guideRna_list_bowtie(bowtie_fp):
    # takes filepath of bowtie output from off-target analysis
    bowtie_df = pd.read_csv(bowtie_fp, sep="\t")
    grouped = bowtie_df.groupby("Name")
    guideRna_list = []

    for grp in grouped:
        # only make a UniqueGuideRNA object if gRNA only has one target site
        # in the alignment
        # Column pos in bowtie output is the start of the alignment
        if len(grp[1]) == 1:
            df = grp[1].iloc[0]
            guideRna = UniqueGuideRNA(
                grp[0], df["sequence"], str(df["chrom"]).lstrip("chr"), df["pos"]
            )

            guideRna_list.append(guideRna)
    if len(guideRna_list) == 0:
        logging.warning("No uniquely binding guide rna found!")
    else:
        logging.info(
            "Uniquely binding guide rna : %s"
            % " ".join([g.name for g in guideRna_list])
        )

    return guideRna_list


def create_snp_df(column_names, guide_rna_list, species, archive_url):
    """
    WARNING : tightly coupled to ensembl_snp_query.R
    pos in the
    R snp search cmd should write to a true tmp file to allow for true multi-processing.
    Although if you try to qsub multiple commands using same project directory the whole thing
    will mess up anyway
    """
    from subprocess import PIPE, STDOUT, Popen

    print("Checking for SNPs...\n")
    if len(guide_rna_list) == 0:
        logging.info("No guide_rnas passed into create_snp_df")
        return pd.DataFrame(columns=column_names)

    guide_rna_string = ":".join([g.info() for g in guide_rna_list])
    R_output_csv = "grna_snps.csv.tmp"

    R_call_opts = "--archive {arch} --species {sp} --Guide {grna}".format(
        arch=archive_url, sp=species, grna=guide_rna_string
    )

    snp_call = str(os.path.join(BIN_DIR, "ensembl_snp_query.R ")) + R_call_opts

    process = Popen(snp_call.split(), stdout=PIPE, stderr=PIPE)
    # process = Popen(shlex.split(snp_call.split()), stdout = PIPE, stderr=PIPE)
    process.wait()  # needed to set process.returncode below
    snp_call_output = process.stdout.read()
    snp_call_error = process.stderr.read()

    logging.debug("snp search R call : %s" % snp_call)
    logging.debug("snp search R stdout : %s" % snp_call_output)
    logging.debug("snp search R stderr : %s" % snp_call_error)

    if process.returncode == 10:
        logging.info("Species %s is not supported for snp searching" % species)
        return pd.DataFrame(columns=column_names)

    try:
        raw_snp_df = pd.read_csv(R_output_csv, sep=",")
    except IOError:
        logging.warning("snp search R call did not work, proceeding without snps")
        return pd.DataFrame()

    os.remove(R_output_csv)

    # combine all snps into one string along with count of snps
    snp_count_list = []
    snp_info_list = []

    # do not display grna name in snp string
    collist = raw_snp_df.columns.tolist()
    collist.remove("Name")

    # reduce the precision of minor_allele_freq and cast to str
    try:
        raw_snp_df["minor_allele_freq"] = raw_snp_df["minor_allele_freq"].apply(
            lambda x: "{0:.3}".format(x)
        )
    except KeyError:
        pass

    for guide_rna in guide_rna_list:

        sub_df = raw_snp_df.loc[raw_snp_df["Name"] == guide_rna.name, collist]
        count = len(sub_df)
        snp_count_list.append(count)

        if count == 0:
            snp_info_list.append(np.NAN)
        else:
            snp_info_list.append(
                "||".join(sub_df.apply(lambda x: ":".join(map(str, x)), axis=1))
            )

    snp_df = pd.DataFrame(
        {column_names[0]: snp_info_list, column_names[1]: snp_count_list},
        index=[g.name for g in guide_rna_list],
    )

    return snp_df


def create_target_bp_df(seq_targeted, guide_rna_list):
    """
    If Distance_from_BP is positive than target bp is UPSTREAM of cas9 cut site
    If negative than target bp is DOWNSTREAM of cas9 cut site

    TODO: refactor so column_names is an input parameter
    """
    # Short-circuit return to all NA df if there are no targeted
    # bases in the seq_targeted
    if seq_targeted.target_indexes:
        target_bp_str = ":".join(map(str, seq_targeted.target_indexes))
    else:
        return pd.DataFrame(
            columns=["BP", "Distance_from_BP"], index=[g.name for g in guide_rna_list]
        )

    distance_from_target = []

    # Only consider guide_rnas that have one target site within the passed
    # sequence, return NAN if this condition is not met
    for guide in guide_rna_list:
        guide_target_list = guide.find_target(seq_targeted)
        # guide rna has no targets on seq_targeted or more than one
        if len(guide_target_list) != 1:
            distance_from_target.append(np.NAN)
            continue

        guide_target = guide_target_list[0]

        raw_distance_from_target = [
            (guide_target.cut - t) for t in seq_targeted.target_indexes
        ]

        adjusted_dist_from_target = [
            d - 1 if d > 0 else d for d in raw_distance_from_target
        ]

        distance_from_target.append(":".join(map(str, adjusted_dist_from_target)))

    return pd.DataFrame(
        {"BP": target_bp_str, "Distance_from_BP": distance_from_target},
        index=[g.name for g in guide_rna_list],
    )


class EntrezGeneClient(object):
    """
    A gene_token is a string representing either a gene name or a integer gene ID
    Be aware that multiple entrez calls in succession can cause connection to
    server to fail.  Separate multiple calls with time.sleep(0.5 to 1)
    """

    def gene_token_to_name(self, gene_token):
        """
        If a name (ie non-int string) is passed it is returned unchanged
        If an id (ie int or int-like string) is passed the resulting
        gene name is upper-cased version of NCBI returned gene name
        """
        token_type = self._parse_gene_token_type(gene_token)

        if token_type == "ID":
            gene_name = self.gene_ID_to_name(gene_token)

            return gene_name

        elif token_type == "NAME":
            return gene_token

        else:
            assert False, "unknown gene_token type %s" % token_type

    def gene_token_to_single_ID(self, gene_token, species):

        token_type = self._parse_gene_token_type(gene_token)

        if token_type == "ID":
            return gene_token

        elif token_type == "NAME":
            gene_id_list = self.gene_name_to_ID(gene_token, species)

            time.sleep(1)

            if len(gene_id_list) > 1:
                message = (
                    "Cmdline gene argument {0} " + "has multiple geneID matches: {1}"
                ).format(gene_token, " ".join(gene_id_list))

                print("{} ".format(message))
                try:
                    picked_id = input("Enter desired geneID -> ")
                    if picked_id in gene_id_list:
                        return picked_id
                    else:
                        print("That geneID was not in the list")
                        assert False
                except EOFError:
                    return None

            # assert False, message

            elif len(gene_id_list) == 0:
                message = (
                    "Cmdline gene argument {0} " + "has no geneID matches"
                ).format(gene_token)

                assert False, message

            else:
                return gene_id_list[0]

        else:
            assert False, "unknown gene_token type %s" % token_type

    def fetch_gbk_record(self, gene_id, extend_length=0):
        """
        gene_id is an integer
        extend_length is the number of bases to include upstream AND downstream
        of the gene region as defined by ncbi
        """
        gene_id = int(gene_id)
        logging.info("NCBI Gene ID : %s" % gene_id)
        gene_summary = self.gene_ID_to_summary(gene_id)

        # pdb.set_trace()
        # genomic_info = gene_summary[0]['GenomicInfo']
        genomic_info = gene_summary["GenomicInfo"]
        assert len(genomic_info) == 1, "more than one chromosomal genbank record"
        # ID for the primary reference chromosome the gene is located on
        genomic_record_id = genomic_info[0]["ChrAccVer"]

        # for some reason these coords are different by one than what you get fromy
        # manually searching ncbi gene database
        seq_start = int(genomic_info[0]["ChrStart"]) + 1
        seq_stop = int(genomic_info[0]["ChrStop"]) + 1

        strand = "1" if seq_start < seq_stop else "2"

        # Biopython freaks out if seq_stop < seq_start
        seq_start, seq_stop = sorted([seq_start, seq_stop])

        seq_start -= int(extend_length)
        seq_stop += int(extend_length)

        entrez_keywords = {
            "id": genomic_record_id,
            "rettype": "gbwithparts",
            "retmode": "text",
            "seq_start": seq_start,
            "seq_stop": seq_stop,
            "strand": strand,
        }

        logging.info("Fetching genbank file with kwargs : %s" % str(entrez_keywords))

        gbk_entrez_handle = Entrez.efetch(db="nuccore", **entrez_keywords)

        gbk_record = SeqIO.read(gbk_entrez_handle, format="genbank")
        gbk_entrez_handle.close()

        return gbk_record

    def gene_name_to_ID(self, gene_name, species):
        """
        Returns a list of all matching gene IDs
        """
        esearch_term = "{sp}[Orgn] AND {gene}[Gene] AND alive[property]".format(
            sp=species, gene=gene_name
        )

        gene_id_list = Entrez.read(
            Entrez.esearch("gene", term=esearch_term, retmax=20)
        )["IdList"]

        return gene_id_list

    def gene_ID_to_summary(self, geneID):
        """
        MAIN FAILURE POINT OF ENTIRE PROGRAM. NCBI XML STYLE
        CHANGES WILL CAUSE THIS TO START THROWING KEYERRORS / INDEXERRORS
        """
        geneID = int(geneID)
        assert geneID >= 0, "Negative gene IDs not allowed : %s" % geneID
        try:
            rec = Entrez.read(Entrez.esummary(id=geneID, db="gene"))
        except RuntimeError:
            raise Exception("NCBI gene ID %s does not exist" % geneID)

        # strip away the style of the raw record and just return the dictionary
        # that has the actual gene info
        return rec["DocumentSummarySet"]["DocumentSummary"][0]

    def gene_ID_to_name(self, geneID):
        return self.gene_ID_to_summary(geneID)["Name"].upper()

    def _filter_old_ID(self, gene_id_list):
        """
        Depreciated in favor of alive[property]
        """
        summary_rec_list = [(id, self.gene_ID_to_summary(id)) for id in gene_id_list]

        currentID_list = [(id, rec[0]["CurrentID"]) for id, rec in summary_rec_list]

        up_to_date_ID = [
            id for id, currentID in currentID_list if str(currentID) == "0"
        ]

        return up_to_date_ID

    def _parse_gene_token_type(self, gene_token):
        try:
            int(gene_token)
            return "ID"

        except ValueError:
            return "NAME"


def find_gRNAs(your_Sequence, your_Gene, experiment):
    # ensure any sequence objects are converted to strings
    # and then converted to upper case
    your_Sequence = str(your_Sequence).upper()
    all_gRNAs = []
    all_gRNAs_namer = []
    counter = 1
    ortho = pd.read_csv(
        "/research_jude/rgs01_jude/groups/millergrp/projects/CrisprDesigns/common/crispr_designs/ortholog12a.csv"
    )
    for g in range(len(ortho.gRNA)):
        print(ortho.gRNA[g])
        print(("Searching for gRNAs to %s" % your_Gene))
        print(("\nYour sequence- > %s" % your_Sequence))
        reverse_seq = Seq(
            your_Sequence
        )  # Create a Biopython Seq object from your_Sequence
        reverse_seq = (
            reverse_seq.reverse_complement()
        )  # Create the reverse_complement of your_Sequence
        revstr = str(reverse_seq)  # make string of reverse_seq
        forward_gRNAs = SeqUtils.nt_search(
            your_Sequence, ortho.gRNA[g]
        )  # search forward seq for gRNAs
        reverse_gRNAs = SeqUtils.nt_search(
            revstr, ortho.gRNA[g]
        )  # search reverse sequence for gRNAs
        forward_gRNAs = forward_gRNAs[1::]  # get rid of gRNA search seq 0index
        reverse_gRNAs = reverse_gRNAs[1::]  # get rid of gRNA search seq 0index
        for x in forward_gRNAs:
            Fwd_Beg = x + ortho.forwardBegin[g]
            Fwd_End = x + ortho.forwardEnd[g]
            PrimeSeq3 = ortho.PAM3[g]
            if PrimeSeq3 is np.NAN:
                PrimeSeq3 = ""
            PrimeSeq5 = ortho.PAM5[g]
            if PrimeSeq5 is np.NAN:
                PrimeSeq5 = ""
            all_gRNAs.append(
                str(PrimeSeq5) + your_Sequence[Fwd_Beg:Fwd_End] + str(PrimeSeq3)
            )

        for x in reverse_gRNAs:
            Fwd_Beg = x + ortho.forwardBegin[g]
            Fwd_End = x + ortho.forwardEnd[g]
            PrimeSeq3 = ortho.PAM3[g]
            if PrimeSeq3 is np.NAN:
                PrimeSeq3 = ""
            PrimeSeq5 = ortho.PAM5[g]
            if PrimeSeq5 is np.NAN:
                PrimeSeq5 = ""
            all_gRNAs.append(str(PrimeSeq5) + revstr[Fwd_Beg:Fwd_End] + str(PrimeSeq3))
    # all_gRNAs = [x for x in all_gRNAs if "TTTT" not in x]   #get rid of gRNAs that have "TTTT"   date 042417
    # print 'gRNAs found that do not have a TTTT tract..'

    #Names the guides
    for i, gRNA in enumerate(all_gRNAs):
        print(i, g)
        if gRNA[0:4] == "TTTV":
            name = your_Gene + str(".Cas12a.g") + str(i + 1)
        elif gRNA[20:23] == "NGG":
            name = your_Gene + str(".Cas9.g") + str(i + 1)
        elif gRNA[21:22] == "CG" or "AG" or "TG":
            name = your_Gene + str(".Cas9NG.g") + str(i + 1)
        all_gRNAs_namer.append(
            name
        )  # create a list that has gRNA names in the format:  your_Gene.gXX
        print("%s %s" % (name, gRNA))

    all_gRNAs_full_name = [experiment + ".%s" % x for x in all_gRNAs_namer]

    return all_gRNAs, all_gRNAs_full_name

def write_CRISPR_Fa(all_gRNAs, all_gRNAs_namer):
    """Create the CRISPR.fa file that will be used for OTA"""
    with open("CRISPR.fa", "w") as f:
        counter = 0
        for x in all_gRNAs_namer:
            f.write(f">{x}\n")  # Guide name
            f.write(f"{all_gRNAs[counter]}\n")  # Guide sequence
            counter += 1
def write_CRISPR_Fa(
    all_gRNAs, all_gRNAs_namer
):  # Create the CRISPR.fa file that will be used for OTA
    counter = 0
    f = open("CRISPR.fa", "w")
    for x in all_gRNAs_namer:
        fa_name = ">" + x
        # print(x, all_gRNAs[counter])
        f.write(fa_name)
        f.write("\n")
        f.write(all_gRNAs[counter])
        f.write("\n")
        print(fa_name)
        print((all_gRNAs[counter]))
        counter = counter + 1
    f.close()

def make_oligos(your_Gene, all_gRNAs, all_gRNAs_namer, vector):
    """
    Makes the oligos that will need to be ordered to clone gRNAs
    """

    # counter = 0
    # mylist = []
    print("Designing oligos...")
    tsv_temp_name = your_Gene + "_oligos.tsv"
    csv_temp_name = your_Gene + "_oligos.csv"

    v = {
        "name": vector.upper(),
        "top_oligo_head": "",
        "top_oligo_tail": "",
        "bottom_oligo_head": "",
        "bottom_oligo_tail": "",
    }
    name = v["name"]

    if name in ("SM115", "SNP23"):
        v["top_oligo_head"] = "ACACCG"
        v["top_oligo_tail"] = "G"
        v["bottom_oligo_head"] = "AAAAC"
        v["bottom_oligo_tail"] = "CG"
    elif name in ("SM102", "SNP21"):
        v["top_oligo_head"] = "TAGG"
        v["top_oligo_tail"] = ""
        v["bottom_oligo_head"] = "AAAC"
        v["bottom_oligo_tail"] = ""
        # v["top_oligo_head_G"] = "TAG"
        # v["top_oligo_head_GG"] = "TA"
    elif name in (
        "PC291",
        "PC323",
        "SM104",
        "SM388",
        "SM428",
        "SNP22",
        "SNP27",
        "SNP28",
        "SS32",
        "SNP38",
        "SP78",
        "SYN",
    ):
        v["top_oligo_head"] = "CACCG"
        v["top_oligo_tail"] = ""
        v["bottom_oligo_head"] = "AAAC"
        v["bottom_oligo_tail"] = "C"
    elif name in ("SM387", "SP82", "SNP36", "SNP112", "SP16", "SNP329"):
        v["top_oligo_head"] = "CACCG"
        v["top_oligo_tail"] = ""
        v["bottom_oligo_head"] = "AAAC"
        v["bottom_oligo_tail"] = "C"
    elif name == "SP76":
        v["top_oligo_head"] = "GAAACACCG"
        v["top_oligo_tail"] = "GTTT"
        v["bottom_oligo_head"] = "TCTAAAAC"
        v["bottom_oligo_tail"] = "CGGTG"
    else:
        assert False, f"Invalid vector choice: {vector}"

    def _make_oligo_file(fp):
        counter = 0
        with open(fp, "w") as f:
            # write header to the oligo output file
            headers = [
                "Name",
                "grna",
                "fwd_oligo_name",
                "fwd_oligo",
                "rev_oligo_name",
                "rev_oligo",
            ]
            headerLine = "\t".join(headers) + "\n"
            f.write(headerLine)

            for x in all_gRNAs:
                mylist = []
                gRNA = x[0:20]
                gRNA = Seq(gRNA)
                if gRNA[0] in "ATCG":
                    top_oligo = v["top_oligo_head"] + str(gRNA) + v["top_oligo_tail"]
                    bottom_oligo = (
                        v["bottom_oligo_head"]
                        + gRNA.reverse_complement()
                        + v["bottom_oligo_tail"]
                    )
                else:
                    print("gRNA does not start with a nucleotide")
                temp_namer = all_gRNAs_namer[counter]
                oligo_name_top = temp_namer + ".F"
                oligo_name_bottom = temp_namer + ".R"
                mylist.extend(
                    [
                        str(all_gRNAs_namer[counter]),
                        all_gRNAs[counter],
                        str(oligo_name_top),
                        str(top_oligo),
                        str(oligo_name_bottom),
                        str(bottom_oligo),
                    ]
                )

                f.write("\t".join(mylist) + "\n")

                # newLineString = ""
                # for s in [
                #         all_gRNAs_namer[counter],
                #         all_gRNAs[counter],
                #         oligo_name_top,
                #         top_oligo,
                #         oligo_name_bottom,
                #         bottom_oligo,
                #     ]:
                #     if s != 0:
                #         newLineString += "\t"
                #     newLineString += f"{s}"
                # f.write(newLineString + "\n")

                counter += 1

    _make_oligo_file(tsv_temp_name)
    _make_oligo_file(csv_temp_name)
    return tsv_temp_name


def make_project_dir(**kwargs):
    options = {}
    options.update(kwargs)
    sep = "-"
    abbrv = SPECIES_TO_ABBRV[options["species"].lower()]
    results = make_unique_file(
        sep.join(
            [
                abbrv + options["gene"],
                options["customer"],
                options["experiment"].rstrip(sep),
            ]
        )
    )
    # print(f"*** Make project dir results: {results}")
    return results


def make_unique_file(filename_template):
    fp = os.path.join(DESIGNS_DIR, filename_template)
    i = 1
    while os.path.exists(fp):
        fp = os.path.join(DESIGNS_DIR, filename_template + "." + str(i))
        i += 1
    return fp


def mark_grna_as_selected(selected_guide_rna, summary_fp):
    df = read_crispr_summary(summary_fp)

    df["Picked"].update(pd.Series(True, index=[g.name for g in selected_guide_rna]))

    summary_write(summary_fp, df)
    return None


def retrieve_features_for_genbank(sequence_column, summary_fp):
    """
    Primer columns are hardcoded.
    Returns a list of 3-tuples (qualifier_dict, feature_type, Bio.Seq)
    """
    df = read_crispr_summary(summary_fp)
    df = df.loc[df["Picked"]]
    grna_features = [
        ({"note": g}, "misc_feature", Seq(s))
        for g, s in zip(df.index, df[sequence_column])
    ]

    # select only rows with amplicons
    # Likely to have redundant entries in dataframe, initially only take seq as a string
    # to facilitate hassle-free hashing into set.  Finally construct Bio.Seq.Seq objects
    df = df.loc[df["Amplicon_ID"].notnull()]
    fwd_primer_features = [
        (p, "primer_bind", s) for p, s in zip(df["Fwd_Primer_Name"], df["Fwd_Primer"])
    ]

    rev_primer_features = [
        (p, "primer_bind", s) for p, s in zip(df["Rev_Primer_Name"], df["Rev_Primer"])
    ]

    primer_features = [
        ({"note": p}, t, Seq(s))
        for p, t, s in set(fwd_primer_features + rev_primer_features)
    ]

    return primer_features + grna_features


def safe_remove(fp):
    try:
        os.remove(fp)
    except OSError:
        pass


def select_grna_from_substring_list(substr_list, guide_rna_list):
    substr_id = lambda x: x.split(".")[-1]
    substr_sorted = sorted(substr_list, key=substr_id)

    selected_grna = []

    for id, grp in itertools.groupby(substr_sorted, key=substr_id):

        matches = [g for g in guide_rna_list if g.get_relative_id() == id]

        selected_grna.extend(matches)

        assert len(matches) > 0, "No guide RNA matches for : %s" % " ".join(grp)

    return selected_grna


def summary_write(output_fp, *args, **kwargs):
    """
    Preserves index ordering of first dataframe if no columns kwarg passed
    """
    index_order = args[0].index

    final_df = pd.concat(args, sort=True, axis=1, join="outer")
    final_df.dropna(subset=["gRNA"], inplace=True)  # Removes gRNA duplicates
    # final_df.fillna(
    #     {"fwd_oligo_name": "", "fwd_oligo": "", "rev_oligo_name": "", "rev_oligo": "",},
    #     inplace=True,
    # )
    final_df = final_df.reindex(index_order, axis=0)  # Reset row numbers w/o dupes

    # Uncomment block below to allow sorting on kwargs
    # if len(kwargs) > 0:
    #   kwargs["inplace"] = True
    #   final_df.sort_values(**kwargs)

    final_df = final_df.astype(
        {
            "long_0": "Int64",
            "long_1": "Int64",
            "long_2": "Int64",
            "short_0": "Int64",
            "fwd_oligo_name": "str",
            "fwd_oligo": "str",
            "rev_oligo_name": "str",
            "rev_oligo": "str",
        }
    )
    final_df.sort_values(by=["long_0", "long_1", "long_2", "short_0"], inplace=True)
    final_df.to_csv(output_fp, sep="\t", na_rep="NA", index_label="Name")
    final_df_temp = pd.DataFrame()
    final_df_temp = final_df
    try:
        final_df_temp = final_df_temp.drop(["BP"], axis=1)  # Removed: 'SNP_count',
        final_df_temp.to_csv(
            "SUMMARY_easy.txt", sep="\t", na_rep="NA", index_label="Name"
        )
    except ValueError:
        print("value error")
        final_df_temp.to_csv(
            "SUMMARY_easy.txt", sep="\t", na_rep="NA", index_label="Name"
        )
        pass
    # print final_df
    return final_df


def copy_folder_over(project_dir):
    try:
        while True:
            choice = input(
                f"Copy new folder ({project_dir}) to CORE PROJECTS? (y/n)\n  -> "
            )
            print()
            if choice.upper() != "Y" and choice.upper() != "N":
                print("not a valid option")
                continue

            elif choice.upper() == "Y":
                source_fp = os.path.join(DESIGNS_DIR, project_dir)
                target_fp = os.path.join(Z_DESIGNS_FOR_CLUSTER, project_dir)
                # print(Z_DESIGNS_FOR_CLUSTER)
                # print(os.path.abspath(Z_DESIGNS_FOR_CLUSTER))
                # time.sleep(3)
                os.system(f"rsync -a {source_fp} {Z_DESIGNS_FOR_CLUSTER}")
                # time.sleep(20)
                # os.system(f"mv {target_fp}_temp {target_fp}")
                # print("Copying over...")
                # time.sleep(10)

                full_windows_dir = PureWindowsPath(Z_DESIGNS_DIR / project_dir)
                print("New folder on the Z drive:")
                # print("    WINDOWS:")
                # print(f"        {full_windows_dir}")
                # print()
                print(f"    explorer {full_windows_dir}")

                return 0
            elif choice.upper() == "N":
                break
    except EOFError:
        return None


def check_primer_off_targets(primer_df, project_directory):
    # Check off targets of primers designed automatically using primer3
    print("\nChecking primers for off target binding sites...")
    primer_df = primer_df.dropna(subset=["Fwd_Primer"])
    primer_dict = dict(zip(primer_df.Fwd_Primer_Name, primer_df.Fwd_Primer))
    primer_dict.update(dict(zip(primer_df.Rev_Primer_Name, primer_df.Rev_Primer)))
    with open("primers.fa", "w") as f:
        for key, value in primer_dict.items():
            f.write(str(">" + key))
            f.write("\n")
            f.write(value)
            f.write("\n")
    off_target_all_primers(
        os.path.join(DESIGNS_DIR, project_directory),
        ABBRV_TO_SPECIES[re.split(r"[A-Z]", project_directory)[0]],
    )


def check_primer_off_targets_dict(primer_dict, project_directory):
    # Check off targets of primers designed automatically using primer3
    print("\nChecking primers for off target binding sites...")
    with open("primers.fa", "w") as f:
        for key, value in primer_dict.items():
            f.write(str(">" + key))
            f.write("\n")
            f.write(value)
            f.write("\n")
    off_target_all_primers(
        os.path.join(DESIGNS_DIR, project_directory),
        ABBRV_TO_SPECIES[re.split(r"[A-Z]", project_directory)[0]],
    )


def read_crispr_summary(fp):
    return pd.read_csv(fp, sep="\t", index_col="Name")


def read_params(fp):
    return pickle.load(open(fp, "rb"))


def remove_all_but_these_files(excluded_fp_list):
    files = glob.glob("*")
    [os.remove(fp) for fp in files if fp not in excluded_fp_list]


def run_command_and_log(cmd):
    print(cmd)
    print(output)

    # round_to_n = lambda x, n: round(x, -int(floor(log10(abs(x)))) + (n - 1))


def _do_formatting(value, format_str):
    """Format value according to format_str, and deal
    sensibly with format_str if it is missing or invalid.
    """
    if format_str == "":
        if type(value) in six.integer_types:
            format_str = ","
        elif type(value) is float:
            format_str = "f"
        elif type(value) is str:
            format_str = "s"
    elif format_str[0] == ".":
        if format_str.endswith("R"):
            if type(value) in six.integer_types:
                value = round_to_n(value, int(format_str[1]))
                format_str = ","
        if not format_str.endswith("G"):
            format_str = format_str + "G"
    try:
        value = format(value, format_str)
    except:
        value = format(value, "")

    return value


def process_position_parameter(param):
    """Process positioning parameters (left, top, width, height) given to
    df_to_table.
    If an integer, returns the right instance of the Cm class to allow it to be
    treated as cm. If missing, then default to 4cm. Otherwise, pass through
    whatever it gets.
    """
    if param is None:
        return Cm(4)
    elif type(param) is int:
        return Cm(param)
    else:
        return param


def df_to_table(
    slide,
    df,
    picked,
    left=Inches(1),
    top=Inches(1.5),
    width=Inches(5.5),
    height=None,
    colnames=None,
    col_formatters=None,
    rounding=None,
    name="CRISPR_summary",
):
    """Converts a Pandas DataFrame to a PowerPoint table on the given
    Slide of a PowerPoint presentation. Look for complete instructions at 
    https://github.com/robintw/PandasToPowerpoint/blob/master/pd2ppt/pd2ppt.py
    """
    left = process_position_parameter(left)
    top = process_position_parameter(top)
    width = process_position_parameter(width)
    height = process_position_parameter(height)

    rows, cols = df.shape
    shp = slide.shapes.add_table(rows + 1, cols, left, top, width, height)

    if colnames is None:
        colnames = list(df.columns)

    # Insert the column names
    for col_index, col_name in enumerate(colnames):
        shp.table.cell(0, col_index).text = col_name
        cell = shp.table.cell(0, col_index)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(9)

    m = df.values

    for row in range(rows):
        for col in range(cols):
            val = m[row, col]

            if col_formatters is None:
                text = str(val)
            else:
                text = _do_formatting(val, col_formatters[col])

            shp.table.cell(row + 1, col).text = text
            cell = shp.table.cell(row + 1, col)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(9)
            para.font.name = "Courier New"
            if row in picked:
                para.font.color.rgb = RGBColor(255, 0, 0)
            shp.table.cell(
                row, col
            ).text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            cell.text_frame.word_wrap = False
    shp.table.columns[0].width = Inches(1.75)
    shp.table.columns[1].width = Inches(2.25)
    ##
    try:
        shp.table.columns[7].width = Inches(1.5)
    except:
        pass
        # print("no column 8")

    if name is not None:
        shp.name = name

    return shp


def CAGE_slides(gene, species, project_dir, edit):
    s_gene = species[0] + gene
    presentation_filename = f"{project_dir}_{edit}_Summary.pptx"
    final_path = os.path.join(DESIGNS_DIR, project_dir, presentation_filename)
    prs = Presentation(os.path.join(BIN_DIR, "empty.pptx"))
    # slide 1
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    shapes.title.text = f"{s_gene} {edit} Summary"

    left = Inches(6)
    top = Inches(1.5)
    width = Inches(0.3)
    height = Inches(0.75)

    shape = shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 0, 0)
    line = shape.line
    line.color.rgb = RGBColor(255, 0, 0)

    # slide 2
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    shapes.title.text = f"{s_gene} gRNA(s)"
    top = width = height = Inches(1.5)
    left = Inches(2.75)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "gRNA sequences indicated with gray arrows"

    # slide 3
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Off Target Analysis - " + s_gene
    left = Inches(2.11)  # Changed from 2.5 on 11/16/2020
    top = Inches(5.45)  # Changed from 4.5 on 11/16/2020
    pic = slide.shapes.add_picture(os.path.join(BIN_DIR, "off_target.jpg"), left, top,)

    try:
        df = pd.read_csv(FINAL_SUMMARY, sep=r"\t", engine="python")
        picked = df[df["Picked"] == True].index.tolist()
        list_df = df.loc[picked, "Name"]
        # gRNA_list = df.loc[picked,'Name'].to_string(index = False)
        if "SNP" in df.columns:
            if isnan(df.loc[0, "Distance_from_BP"]):
                df1 = df.iloc[:, 0:7]
            else:
                df1 = df.iloc[:, np.r_[0:7, 9]]
        else:  # If SNP_count is not there because ENSEMBL was down or smthng
            if isnan(df.loc[0, "Distance_from_BP"]):
                df1 = df.iloc[:, 0:6]
            else:
                df1 = df.iloc[:, np.r_[0:6, 7]]
        df1 = df1.fillna("NA")
        # df1 = df1.head(14)  # To make everything fit above the off_target image
        table = df_to_table(slide, df1, picked)
    except:
        print("Trouble making table in presentation. Please make manually.")

    # slide 4
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Summary"

    tf = body_shape.text_frame
    tf.text = "The following gRNA(s) are recommended based on off-target profile and the parameters of your project:"

    try:
        for i in list_df:
            p = tf.add_paragraph()
            p.text = i
            p.level = 1
    except:
        p = tf.add_paragraph()
        p.text = "add your gRNA names here"
        p.level = 1

    # ppt created
    prs.save(final_path)

    # TODO: Try adding prs.close() OR enclosing above in a "with ... as prs" context

    return final_path


def main():
    parser = argparse.ArgumentParser()
    subparser = parser.add_subparsers()

    # Project initialization command
    init_parser = subparser.add_parser("init")

    init_parser.add_argument(
        "-c", "--customer", required=True, help="name of the customer"
    )

    init_parser.add_argument(
        "-g",
        "--gene",
        required=True,
        help="name or NCBI gene ID of gene crispr is targeting. "
        + "If a gene ID is passed the resulting directory will be named using the gene name.",
    )

    init_parser.add_argument(
        "-e", "--experiment", required=True, help="string to prepend to grna names"
    )

    init_parser.add_argument(
        "-sp",
        "--species",
        required=True,
        choices=SUPPORTED_SPECIES,
        help="species that cripsr is being made for",
    )

    init_parser.add_argument(
        "-gle",
        "--GENBANK_LENGTH_EXTENSION",
        required=True,
        help="the number of bp on either end of the gene of interest in gene.gbk",
    )

    # init_parser.add_argument(
    #     "-v",
    #     "--vector",
    #     required=True,
    #     choices=VECTOR_CHOICES,
    #     help="Species that cripsr is being made for.\n"
    #     + "Values are given in source code of make_oligos function"
    # 'SP78 (use instead of SNP23!): for Mammalian U6'+\
    # 'SYN: Synthetic gRNA'+\
    # 'SNP23 (aka SM115): for Mammalian U6'+\
    # 'SNP21 (aka SM102): for Zebrafish T7'+\
    # 'SNP22 (aka SM104): for pX330 single vector gRNA + Cas9'+\
    # 'SM387: Lenti-CRISPR'+\
    # 'SNP36 (aka SP82):  Lenti-CRISPRv2'+\
    # 'SNP112: (Lenti-CRISPRv2 w eGFP)'+\
    # 'SNP27 (aka SM388): pX458 single vector gRNA +Cas9-T2A-eGFP (AppliedStemCell)' +\
    # 'SNP38 (aka SS32, PC323) pX459 single vector gRNA +Cas9-T2A-Puro'+\
    # 'SNP28 (aka SM438): lentiGuide-Puro  single vector gRNA + EF1a-Puro (Milbrandt)'+\
    # 'SP16: lentiGuide-eGFP, single vector gRNA that expresses eGFP' +\
    # 'SP76 : Sanger vector'+\=
    # )

    init_parser.set_defaults(func=initialize)

    # crispr design command
    grna_parser = subparser.add_parser("grna")
    grna_parser.add_argument(
        "-s",
        "--sequence",
        required=True,
        type=SeqTargeted,
        help="sequence string in which to search for gRNAs. "
        + "Mark targeted base pairs with <>",
    )

    grna_parser.add_argument(
        "-p",
        "--project_dir",
        required=True,
        help="path to project directory to make guide rna for",
    )

    grna_parser.add_argument(
        "-v",
        "--vector",
        type=str,
        help="Values are given in source code of make_oligos function",
    )

    grna_parser.add_argument(
        "--force",
        action="store_true",
        help="instructs program to continue even if "
        + "%s is in target directory" % CRISPR_SUMMARY,
    )

    grna_parser.set_defaults(func=crispr_design)

    # Primer design command
    primer_parser = subparser.add_parser("primer")
    primer_parser.add_argument(
        "-p",
        "--project_dir",
        required=True,
        help="path to project directory to make pcr assay primers for",
    )

    primer_parser.add_argument(
        "-i",
        "--grna_id",
        nargs="+",
        required=True,
        help="space separated list of gRNA ids to make primers for. Crispr ids"
        + " are given by the Name column in %s" % CRISPR_SUMMARY,
    )

    primer_parser.add_argument(
        "--fuzzy",
        help="string detailing what kind of fuzzy match should be performed "
        + "when aligning gRNA seqs back to reference gene.gbk.  Of form [sid]<=\d surrounded by {} and "
        + "separated by , . See documentation of fuzzy matching for python module regex (not re)",
    )

    primer_parser.set_defaults(func=primer_design)

    joiner_parser = subparser.add_parser("srmjoiner")
    joiner_parser.set_defaults(func=SRMJoiner)

    args = parser.parse_args()
    args.func(args)
    return None


if __name__ == "__main__":

    def backup_project_directory(args, params):
        """
        Delete the results of primer command because different gRNAs will be picked
        Also delete previous results of grna command.
        All previous files are backed-up in a *.1, *.2,... directory
        BE CAREFUL USING THIS. DOES AN OS.CHDIR DANCE
        """

        os.chdir("../")
        backup_dir = make_project_dir(**params)  # Adds to .[num] to the end of fp
        shutil.copytree(args.project_dir, backup_dir)

        os.chdir(args.project_dir)
        remove_all_but_these_files([PROJECT_PARAMS, GENE_GENBANK, LOG_FILE])

        # Don't delete log file, just write it empty so logging module
        # will work correctly
        with open(LOG_FILE, "w") as fh:
            fh.write("")

        return backup_dir

    def handle_project_dir_state_for_crispr_design(args):
        # do not continue if grnas have already been picked in the project directory
        if os.path.exists(FINAL_SUMMARY) and not args.force:
            assert False, (
                "assay primers have already been created, call grna command with "
                + "--force flag to continue anyway. Target directory will be copied into "
                + "new directory with a .1, .2, ... appended"
            )

        # If --force flag is passed, args.project directory and then delete
        # all files created after init and wipe log
        elif args.force and os.path.exists(FINAL_SUMMARY):
            params = read_params(PROJECT_PARAMS)
            backup_dir = backup_project_directory(args, params)
            logging.info(
                "Executing grna under --force flag; previous run backed-up "
                + "in directory %s" % backup_dir
            )

        # If no primers made, delete all files not created by initialize() ie do a fresh run
        else:
            remove_all_but_these_files([PROJECT_PARAMS, GENE_GENBANK, LOG_FILE])

        return None

    def crispr_design(args):
        # change working directory to project directory because all global filepath
        # variables are relative to project dir
        os.chdir(os.path.join(DESIGNS_DIR, args.project_dir))

        logging.basicConfig(filename=LOG_FILE, level=logging.DEBUG)

        handle_project_dir_state_for_crispr_design(args)

        # update params pickle-file with the sequence used for gRNA picking
        params = read_params(PROJECT_PARAMS)
        with open(PROJECT_PARAMS, "wb") as fh:
            new_params = params.copy()
            new_params["grna_picking_seq"] = args.sequence
            pickle.dump(new_params, fh, -1)
            del new_params

        logging.info(
            "Sequence selected for crispr design : %s" % args.sequence.raw_sequence
        )
        logging.info("Vector choice : %s" % args.vector)

        all_gRNAs, all_gRNAs_namer = find_gRNAs(
            args.sequence, params["gene"], params["experiment"]
        )

        # Make guideRNA oligos
        oligo_summary_fp = make_oligos(
            params["gene"], all_gRNAs, all_gRNAs_namer, args.vector
        )

        oligo_summary_df = pd.read_csv(oligo_summary_fp, sep="\t", index_col=0)
        oligo_df = oligo_summary_df.loc[
            :, ["fwd_oligo_name", "fwd_oligo", "rev_oligo_name", "rev_oligo"]
        ]

        write_CRISPR_Fa(all_gRNAs, all_gRNAs_namer)

        # This is the same dataframe as represented by
        # CRISPR_off_target_v3.summary
        summary_df = off_target_all(params["species"].lower())
        customer_df = pd.DataFrame(
            {CUSTOMER_COLUMN: params["customer"]}, index=all_gRNAs_namer
        )
        # Create a column that will record whether a grna has been picked by the
        # primer command.
        picked_df = pd.DataFrame({"Picked": False}, index=all_gRNAs_namer)
        guide_rna_list = [
            GuideRNA(name, g) for name, g in zip(all_gRNAs_namer, all_gRNAs)
        ]
        # If Distance_from_BP is positive than target bp is UPSTREAM of cas9 cut site
        # If negative than target bp is DOWNSTREAM of cas9 cut site
        target_bp_df = create_target_bp_df(args.sequence, guide_rna_list)

        # Only creates guide_rna objects for grna with long_0 == 1
        unique_guide_rna_list = create_unique_guideRna_list_bowtie(OFF_TARGET_BOWTIE)
        # calls R subcommand then reads output R dataframe into memory
        # also deletes the temporary R-generated csv
        snp_df = create_snp_df(
            SNP_DATAFRAME_COLUMNS,
            unique_guide_rna_list,
            species=params["species"],
            archive_url=ENSEMBL_RELEASE_ARCHIVE,
        )

        # Fault tolerance against snp search, if snp search fails snp_df
        # will be an empty dataframe
        if snp_df.empty:
            sort_key = BASIC_SUMMARY_SORT_KEY
        else:
            sort_key = BASIC_SUMMARY_SORT_KEY + list(SNP_DATAFRAME_COLUMNS)[1:2]
        # combine all the loose dataframes by columns in the order they are passed
        # also sorts rows by sort_key
        summary_write(
            CRISPR_SUMMARY,
            summary_df,
            snp_df,
            target_bp_df,
            customer_df,
            picked_df,
            oligo_df,
            columns=sort_key,
        )
        merged_summary_df = pd.merge(
            summary_df,
            target_bp_df[["Distance_from_BP"]],
            left_index=True,
            right_index=True,
            how="left",
        )

        # Add SNPs to end of summary, if SNP working
        # merged_summary_df = pd.merge(
        #     merged_summary_df,
        #     snp_df, # Column index 1
        #     left_index=True,
        #     right_index=True,
        #     how='left',
        # )

        print("Summary:")
        print(merged_summary_df)
        # copies and pastes in empty pptx slidedeck
        # os.system('cp /research/rgs01/project_space/millergrp/CrisprDesigns/common/local/crispr_design/template_KO_Summary.pptx '
        #   + args.project_dir.split('-')[0] + '_Summary.pptx')

        # Print path to summary file here?

        return None

    # rom interactive_crispr_NGS_python3 import GENBANK_LENGTH_EXTENSION

    def initialize(args):
        # Handle passing a ncbi geneID as the gene
        entrezClient = EntrezGeneClient()

        gene_name = entrezClient.gene_token_to_name(args.gene)

        # This method will throw error and kill program before making the project directory
        # if passed gene is argument ambiguous (ie multiple matches) or fake (no matches)
        geneID = entrezClient.gene_token_to_single_ID(args.gene, args.species)

        dir_name_values = vars(args).copy()
        dir_name_values["gene"] = gene_name

        project_dir = make_project_dir(**dir_name_values)
        new_dir_on_windows = PureWindowsPath(
            project_dir.replace("/research_jude", "").replace(
                "/rgs01_jude/groups/millergrp/projects/",
                r"Z:\ResearchHome\Groups\millergrp\projects\\",
            )
        )

        assert not os.path.isdir(project_dir), f"{project_dir} already exists!"

        os.mkdir(project_dir)
        os.chdir(project_dir)

        print("\nNew project directory:")
        print(f"    On Cluster:\n        {project_dir}")
        print(f"\n    On Windows:\n        {new_dir_on_windows}")

        logging.basicConfig(filename=LOG_FILE, level=logging.DEBUG)
        logging.info("Initialized project directory %s" % os.getcwd())
        logging.info("Gene Commandline : %s" % args.gene)
        logging.info("Gene Name: %s" % gene_name)
        logging.info("Species : %s" % args.species)
        logging.info("Customer : %s" % args.customer)

        gbk_record = entrezClient.fetch_gbk_record(
            geneID, extend_length=int(args.GENBANK_LENGTH_EXTENSION)
        )

        SeqIO.write([gbk_record], open(GENE_GENBANK, "w"), "genbank")

        check_projects(gene_name, args.species)

        # write the command line parameters as a pickle file so state can be
        # reconstructed by other commands. Substitute cmdline gene
        # for gene_name because params['gene'] used to name things in later commands
        with open(PROJECT_PARAMS, "wb") as fh:
            params = vars(args).copy()
            params.pop("func")
            params["gene"] = gene_name
            pickle.dump(params, fh, -1)

        return None

    def primer_design(args):
        if args.project_dir is not None:
            os.chdir(os.path.join(DESIGNS_DIR, args.project_dir))

        if args.fuzzy is not None and not GuideRNA.validate_fuzzy_match(args.fuzzy):
            raise InvalidFuzzyStringError(args.fuzzy)

        logging.basicConfig(filename=LOG_FILE, level=logging.DEBUG, filemode="a")
        logging.info("Guide-rna selected for primer design %s" % " ".join(args.grna_id))

        params = read_params(PROJECT_PARAMS)
        # Allow all guide RNAs to be selected for primer design
        guide_rna_list = create_guideRna_list_from_summary(
            CRISPR_SUMMARY, SEQUENCE_COLUMN
        )

        selected_guide_rna = select_grna_from_substring_list(
            args.grna_id, guide_rna_list
        )

        # print(selected_guide_rna)

        # Updates the CRISPR_SUMMARY csv with guide rna selections via Picked column
        mark_grna_as_selected(selected_guide_rna, CRISPR_SUMMARY)

        primerDesigner = PrimerDesigner(
            selected_guide_rna,
            SeqIO.read(open(GENE_GENBANK, "r"), format="genbank"),
            params["grna_picking_seq"],
            args.fuzzy,
        )

        ampliconSet = primerDesigner.create_ampliconSet(args=args)
        # print(ampliconSet.amplicon_list)

        amplicon_name = str(
            params["experiment"] + "." + params["gene"]
        )  # give all amplicons the same name
        ampliconSet.write_report_for_N_best_amplicons(
            3, AMPLICON_REPORT, name=amplicon_name
        )

        # chooses the best amplicon and returns a dataframe to be added to CRISPR_summary
        best_amplicon, amplicon_df = ampliconSet.choose_best_amplicon(
            name=amplicon_name
        )
        # print("2207")
        print(
            amplicon_df[
                [
                    "Amplicon_Length",
                    "Fwd_Primer_Name",
                    "Fwd_Primer",
                    "Rev_Primer_Name",
                    "Rev_Primer",
                ]
            ]
        )

        # Write final_summary if not already else update the existing table
        if os.path.exists(FINAL_SUMMARY):
            summary_df = read_crispr_summary(FINAL_SUMMARY)
            # CRISPR_SUMMARY not read so we need to also updated Picked column here
            amplicon_df.insert(0, "Picked", True)
            summary_df.update(amplicon_df)
            temp_df = pd.DataFrame()
            temp_df = summary_write(FINAL_SUMMARY, summary_df)
            # summary_df.loc[1] = SRM_cols
            check_primer_off_targets(temp_df, args.project_dir)
            # summary_write(FINAL_SUMMARY, summary_df)
        else:
            summary_df = read_crispr_summary(CRISPR_SUMMARY)
            # summary_df.loc[1] = SRM_cols
            # print summary_df
            temp_df = pd.DataFrame()
            temp_df = summary_write(FINAL_SUMMARY, summary_df, amplicon_df)
            check_primer_off_targets(temp_df, args.project_dir)
            # summary_write(FINAL_SUMMARY, summary_df, amplicon_df)

        print("\nMaking .XLS file for SRM upload...")
        summarySeries = crisprSummarySRMConverter.CRISPRSummaryToSeries(FINAL_SUMMARY)
        crisprSummarySRMConverter.CRISPRSummarySeriesToXLS(
            summarySeries, "Converted-Picked-Py.xls"
        )

        # finally write grna and amplicon to genbank file
        mod_gene_genbank = params["gene"] + "_mod_NGS.gbk"
        genbank_features = retrieve_features_for_genbank(SEQUENCE_COLUMN, FINAL_SUMMARY)
        genBankWriter = GenBankFeatureWriter(GENE_GENBANK)
        genBankWriter.add_features(genbank_features)
        genBankWriter.write(mod_gene_genbank)
        print("Please enter the type of project to append to PPT Title.")
        print("  Examples: KO, CKO, del, Y38G, tag...")
        print('  ** Please avoid special characters, such as ">" **')
        try:
            project_type = input("\n  -> ")
        except EOFError:
            return None

        ppt_path = CAGE_slides(
            params["gene"], params["species"], args.project_dir, project_type
        )

        print()
        print("*" * 50)
        print("*" + (" " * 48) + "*")
        print(f"* {'! ! !   D E S I G N   C O M P L E T E   ! ! !':^46} *")
        print("*" + (" " * 48) + "*")
        print("*" * 50)
        print()

        copy_folder_over(args.project_dir)
        return None

    def SRMJoiner(args):
        use_default_filename_choice = ""
        filename = ""
        first_folder = ""
        second_folder = ""
        try:
            while use_default_filename_choice not in ["Y", "YES", "N", "NO"]:
                print("Use default filename 'CRISPR_summary_w_primers_NGS.txt'? (y/n)")
                use_default_filename_choice = input("--> ").strip().upper()

            if use_default_filename_choice in ["N", "NO"]:
                print("Please enter the intended filename:")
                filename = input("--> ")
            elif use_default_filename_choice in ["Y", "YES"]:
                filename = "CRISPR_summary_w_primers_NGS.txt"

            # while os.path.isdir(first_folder) == False:
            print("Please enter the full path to the folder of the first file.")
            print("* NOTE: This will serve as the output folder! *")
            first_folder = (
                # '"'
                PureWindowsPath(
                    input("--> ")
                    .strip()
                    .replace(
                        r"Z:\ResearchHome\Groups\millergrp\home\common",
                        r"\research_jude\rgs01_jude\groups\millergrp\home\common",
                    )
                    # .replace("CORE PROJECTS", '"CORE PROJECTS"')
                    # .split("\\")
                    # + '"'
                ).as_posix()
            )

            # while os.path.isdir(second_folder) == False:
            print("Please enter the full path to the folder of the second file.")
            second_folder = (
                # '"'
                PureWindowsPath(
                    input("--> ")
                    .strip()
                    .replace(
                        r"Z:\ResearchHome\Groups\millergrp\home\common",
                        r"\research_jude\rgs01_jude\groups\millergrp\home\common",
                    )
                    # .replace("CORE PROJECTS", '"CORE PROJECTS"')
                    # .split("\\")
                    # + '"'
                ).as_posix()
            )

            firstFile = PurePosixPath(first_folder).joinpath(filename)
            secondFile = PurePosixPath(second_folder).joinpath(filename)
            crisprSummarySRMConverter.dualTXTtoXLS(firstFile, secondFile)
            print("\n")
            return None
        except EOFError:
            print("\n")
            return None

    # RUN THE PROGRAM
    main()
